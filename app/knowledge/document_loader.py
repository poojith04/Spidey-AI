from pathlib import Path

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


class DocumentLoader:

    def load_folder(self, folder):

        documents = []

        folder = Path(folder)

        for file in folder.iterdir():

            if file.suffix.lower() == ".pdf":
                text = self.read_pdf(file)

            elif file.suffix.lower() == ".docx":
                text = self.read_docx(file)

            elif file.suffix.lower() == ".pptx":
                text = self.read_ppt(file)

            elif file.suffix.lower() == ".txt":
                text = self.read_txt(file)

            else:
                continue

            documents.append({
                "file": file.name,
                "text": text
            })

        return documents

    def read_pdf(self, path):

        reader = PdfReader(path)

        text = ""

        for page in reader.pages:
            text += page.extract_text() or ""

        return text

    def read_docx(self, path):

        doc = Document(path)

        return "\n".join(
            paragraph.text
            for paragraph in doc.paragraphs
        )

    def read_ppt(self, path):

        presentation = Presentation(path)

        text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text

    def read_txt(self, path):

        with open(path, "r", encoding="utf-8") as file:
            return file.read()